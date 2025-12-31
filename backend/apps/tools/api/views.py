"""Tools API Views - Complete PDF Processing Endpoints"""
from rest_framework import viewsets, status, permissions
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework.decorators import action
from rest_framework.parsers import MultiPartParser, FormParser
from django.http import HttpResponse, FileResponse
import io
import tempfile
import os
import logging

logger = logging.getLogger(__name__)


class PDFToolAPIView(APIView):
    """Base class for PDF tool endpoints with file upload handling."""
    parser_classes = [MultiPartParser, FormParser]
    permission_classes = [permissions.AllowAny]  # Allow Guests (Limits enforced in check_usage_limit)

    def get_file_from_request(self, request):
        """Extract uploaded file from request."""
        if 'file' not in request.FILES:
            return None, Response({'error': 'No file provided'}, status=status.HTTP_400_BAD_REQUEST)
        return request.FILES['file'], None


# ─────────────────────────────────────────────────────────────────────────────
# CONVERSION TO PDF
# ─────────────────────────────────────────────────────────────────────────────

class WordToPDFView(PDFToolAPIView):
    """Convert Word documents to PDF."""
    
    def post(self, request):
        # Specific check
        allowed, error = check_usage_limit(request, 'WORD_TO_PDF')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from apps.tools.converters.word_to_pdf import convert_word_to_pdf
            result = convert_word_to_pdf(file)
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class ExcelToPDFView(PDFToolAPIView):
    """Convert Excel spreadsheets to PDF."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'EXCEL_TO_PDF')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from apps.tools.converters.office_converter import convert_excel_to_pdf
            result = convert_excel_to_pdf(file)
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PowerpointToPDFView(PDFToolAPIView):
    """Convert PowerPoint presentations to PDF."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PPT_TO_PDF')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from apps.tools.converters.office_converter import convert_powerpoint_to_pdf
            result = convert_powerpoint_to_pdf(file)
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class JPGToPDFView(PDFToolAPIView):
    """Convert JPG images to PDF."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'JPG_TO_PDF')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from PIL import Image
            img = Image.open(file)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            
            output = io.BytesIO()
            img.save(output, format='PDF')
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class HTMLToPDFView(PDFToolAPIView):
    """Convert HTML to PDF."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'HTML_TO_PDF')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from apps.tools.converters.office_converter import convert_html_to_pdf
            result = convert_html_to_pdf(file)
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class MarkdownToPDFView(PDFToolAPIView):
    """Convert Markdown to PDF."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'MARKDOWN_TO_PDF')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from apps.tools.converters.office_converter import convert_markdown_to_pdf
            result = convert_markdown_to_pdf(file)
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# ─────────────────────────────────────────────────────────────────────────────
# CONVERSION FROM PDF
# ─────────────────────────────────────────────────────────────────────────────

class PDFToJPGView(PDFToolAPIView):
    """Convert PDF to JPG images."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PDF_TO_JPG')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            # Convert first page to image
            page = doc.load_page(0)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            
            output = io.BytesIO()
            output.write(pix.tobytes("jpeg"))
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='image/jpeg')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.jpg"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PDFToWordView(PDFToolAPIView):
    """Convert PDF to Word document."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PDF_TO_WORD')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from pdf2docx import Converter
            
            # Write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                tmp_in.write(file.read())
                input_path = tmp_in.name
            
            output_path = input_path.replace('.pdf', '.docx')
            
            # Convert
            cv = Converter(input_path)
            cv.convert(output_path)
            cv.close()
            
            # Read output
            with open(output_path, 'rb') as f:
                output_content = f.read()
            
            # Cleanup
            os.unlink(input_path)
            os.unlink(output_path)
            
            response = HttpResponse(
                output_content,
                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.docx"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PDFToExcelView(PDFToolAPIView):
    """Convert PDF to Excel."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PDF_TO_EXCEL')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        merge_sheets = request.data.get('merge_sheets', 'false').lower() == 'true'
        output_format = request.data.get('output_format', 'xlsx').lower()
        
        try:
            import pdfplumber
            import pandas as pd
            
            # Write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                tmp_in.write(file.read())
                input_path = tmp_in.name
            
            output_ext = '.csv' if output_format == 'csv' else '.xlsx'
            output_path = input_path.replace('.pdf', output_ext)
            
            all_tables = []
            
            with pdfplumber.open(input_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    for j, table in enumerate(tables):
                        if not table: continue
                        
                        # Basic header detection: Use first row
                        if len(table) > 1:
                            # Clean headers to avoid duplicates or empty
                            headers = table[0]
                            # Ensure unique headers
                            headers = [str(h) if h else f"Col_{k}" for k, h in enumerate(headers)]
                            df = pd.DataFrame(table[1:], columns=headers)
                        else:
                            df = pd.DataFrame(table)
                            
                        all_tables.append((f"P{i+1}_T{j+1}", df))
            
            if not all_tables:
                 # Create a dummy dataframe if no tables found
                 all_tables.append(("Sheet1", pd.DataFrame([["No tables detected"]], columns=["Message"])))

            if output_format == 'csv':
                # For CSV, merge all tables
                final_df = pd.concat([t[1] for t in all_tables], ignore_index=True)
                final_df.to_csv(output_path, index=False)
                content_type = 'text/csv'
            else:
                # XLSX
                if merge_sheets:
                    final_df = pd.concat([t[1] for t in all_tables], ignore_index=True)
                    final_df.to_excel(output_path, index=False)
                else:
                    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                        for name, df in all_tables:
                            # Excel sheet names max 31 chars
                            safe_name = name[:31]
                            df.to_excel(writer, sheet_name=safe_name, index=False)
                content_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            
            # Read output
            with open(output_path, 'rb') as f:
                output_content = f.read()
            
            # Cleanup
            os.unlink(input_path)
            os.unlink(output_path)
            
            response = HttpResponse(output_content, content_type=content_type)
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}{output_ext}"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PDFToPowerpointView(PDFToolAPIView):
    """Convert PDF to PowerPoint."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PDF_TO_PPT')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            from pptx import Presentation
            from pptx.util import Inches
            
            # Write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                tmp_in.write(file.read())
                input_path = tmp_in.name
            
            output_path = input_path.replace('.pdf', '.pptx')
            
            # Open PDF
            doc = fitz.open(input_path)
            
            # Create PowerPoint
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for page in doc:
                # Convert page to image
                pix = page.get_pixmap(dpi=150)
                img_path = f"{input_path}_page_{page.number}.png"
                pix.save(img_path)
                
                # Add slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add image to slide
                left = Inches(0)
                top = Inches(0)
                slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
                
                # Cleanup temp image
                os.unlink(img_path)
            
            doc.close()
            
            # Save PowerPoint
            prs.save(output_path)
            
            # Read output
            with open(output_path, 'rb') as f:
                output_content = f.read()
            
            # Cleanup
            os.unlink(input_path)
            os.unlink(output_path)
            
            response = HttpResponse(
                output_content,
                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.pptx"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PDFToHTMLView(PDFToolAPIView):
    """Convert PDF to HTML."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PDF_TO_HTML')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            html_content = ['<!DOCTYPE html><html><head><meta charset="utf-8"><title>Converted PDF</title></head><body>']
            
            for page in doc:
                page_html = page.get_text('html')
                html_content.append(f'<div class="page" style="page-break-after: always;">{page_html}</div>')
            
            html_content.append('</body></html>')
            doc.close()
            
            full_html = '\n'.join(html_content)
            
            response = HttpResponse(full_html, content_type='text/html')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}.html"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PDFToPDFAView(PDFToolAPIView):
    """Convert PDF to PDF/A format."""
    
    def post(self, request):
        allowed, error = check_usage_limit(request, 'PDF_TO_PDFA')
        if not allowed: return error

        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import subprocess
            
            # Write to temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                tmp_in.write(file.read())
                input_path = tmp_in.name
            
            output_path = input_path.replace('.pdf', '_pdfa.pdf')
            
            # Convert using Ghostscript
            result = subprocess.run([
                'gs', '-dPDFA=2', '-dBATCH', '-dNOPAUSE',
                '-sColorConversionStrategy=UseDeviceIndependentColor',
                f'-sOutputFile={output_path}',
                '-sDEVICE=pdfwrite',
                '-dPDFACompatibilityPolicy=1',
                input_path
            ], capture_output=True, text=True, timeout=300)
            
            if result.returncode != 0:
                os.unlink(input_path)
                return Response({'error': f'Conversion failed: {result.stderr}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            
            # Read output
            with open(output_path, 'rb') as f:
                output_content = f.read()
            
            # Cleanup
            os.unlink(input_path)
            os.unlink(output_path)
            
            response = HttpResponse(output_content, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_pdfa.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)



# ─────────────────────────────────────────────────────────────────────────────
# PDF OPTIMIZATION
# ─────────────────────────────────────────────────────────────────────────────


# ─────────────────────────────────────────────────────────────────────────────
# UTILS
# ─────────────────────────────────────────────────────────────────────────────

def get_client_ip(request):
    x_forwarded_for = request.META.get('HTTP_X_FORWARDED_FOR')
    if x_forwarded_for:
        ip = x_forwarded_for.split(',')[0]
    else:
        ip = request.META.get('REMOTE_ADDR')
    return ip

def check_usage_limit(request, feature_code):
    """
    Checks if user has reached limit for feature.
    Supports usage tracking for both Authenticated Users and Guests (via IP).
    Returns (allowed: bool, response: Response | None)
    """
    try:
        from apps.subscriptions.models.subscription import Feature, UserFeatureUsage, UserFeatureOverride, GuestFeatureUsage
        from django.utils import timezone
        
        # Get Feature
        feature = Feature.objects.filter(code=feature_code).first()
        if not feature:
            return True, None # Feature not restricted or unknown

        # 1. GUEST ACCESS
        if not request.user.is_authenticated:
            # Check Premium-Only Restriction for Guests
            if feature.is_premium_default:
                return False, Response(
                    {'error': 'This is a premium feature. Please login or upgrade.'}, 
                    status=status.HTTP_403_FORBIDDEN
                )
            
            # Check Host/IP Limit
            ip = get_client_ip(request)
            if not ip:
                 return False, Response({'error': 'Cannot identify client.'}, status=status.HTTP_400_BAD_REQUEST)

            today = timezone.now().date()
            usage, created = GuestFeatureUsage.objects.get_or_create(ip_address=ip, feature=feature, date=today)
            
            limit = feature.free_limit
            
            if limit > 0 and usage.count >= limit:
                 return False, Response(
                     {'error': f'Daily guest limit of {limit} reached. Please sign up for more.'}, 
                     status=status.HTTP_403_FORBIDDEN
                 )
            
            usage.count += 1
            usage.save()
            return True, None

        # 2. LOGGED-IN USER ACCESS
        user = request.user
        
        # Check Override
        override = UserFeatureOverride.objects.filter(user=user, feature=feature).first()
        if override:
            if override.is_enabled:
                return True, None
            else:
                return False, Response({'error': 'Feature disabled for your account'}, status=status.HTTP_403_FORBIDDEN)
        
        # Check Subscription Tier (Premium Bypass)
        is_premium = False
        if hasattr(user, 'subscription') and user.subscription:
             pass
        
        # If user model has subscription_tier (e.g. from property)
        if getattr(user, 'subscription_tier', 'FREE') in ['PRO', 'ENTERPRISE', 'PREMIUM']:
             return True, None # Unlimited for paid

        # If User is Free and Feature is Premium Only
        if feature.is_premium_default:
             return False, Response(
                 {'error': 'This feature requires a premium subscription.'}, 
                 status=status.HTTP_403_FORBIDDEN
             )

        # Check Free User Limits
        today = timezone.now().date()
        usage, created = UserFeatureUsage.objects.get_or_create(user=user, feature=feature, date=today)
        
        limit = feature.free_limit
        if limit > 0:
            if usage.count >= limit:
                return False, Response(
                    {'error': f'Daily limit of {limit} reached. Upgrade for unlimited access.'}, 
                    status=status.HTTP_403_FORBIDDEN
                )

        # Increment
        usage.count += 1
        usage.save()
        return True, None
        
    except Exception as e:
        print(f"Limit Check Error: {e}")
        return True, None # Fail safe open

class MergePDFView(PDFToolAPIView):
    """Merge multiple PDFs into one."""
    
    def post(self, request):
        # Usage Check
        allowed, error_response = check_usage_limit(request, 'MERGE_PDF')
        if not allowed:
            return error_response

        files = request.FILES.getlist('files')
        if not files or len(files) < 2:
            return Response({'error': 'At least 2 files required for merge'}, status=status.HTTP_400_BAD_REQUEST)

        
        try:
            import fitz
            merged = fitz.open()
            
            for file in files:
                pdf = fitz.open(stream=file.read(), filetype="pdf")
                merged.insert_pdf(pdf)
                pdf.close()
            
            output = io.BytesIO()
            merged.save(output)
            merged.close()
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="merged.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class SplitPDFView(PDFToolAPIView):
    """Split PDF into multiple files."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            import json
            import zipfile
            
            selected_pages = json.loads(request.data.get('selectedPages', '[]'))
            split_mode = request.data.get('splitMode', 'extract')
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            if split_mode == 'extract' and selected_pages:
                # Extract selected pages
                new_doc = fitz.open()
                for page_num in selected_pages:
                    if 0 <= page_num - 1 < len(doc):
                        new_doc.insert_pdf(doc, from_page=page_num-1, to_page=page_num-1)
                
                output = io.BytesIO()
                new_doc.save(output)
                new_doc.close()
                output.seek(0)
                
                response = HttpResponse(output.read(), content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="split.pdf"'
                return response
            else:
                # Split all pages into ZIP
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                    for i in range(len(doc)):
                        page_doc = fitz.open()
                        page_doc.insert_pdf(doc, from_page=i, to_page=i)
                        page_buffer = io.BytesIO()
                        page_doc.save(page_buffer)
                        page_doc.close()
                        page_buffer.seek(0)
                        zip_file.writestr(f'page_{i+1}.pdf', page_buffer.read())
                
                doc.close()
                zip_buffer.seek(0)
                
                response = HttpResponse(zip_buffer.read(), content_type='application/zip')
                response['Content-Disposition'] = 'attachment; filename="split_pages.zip"'
                return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class PDFPagePreviewsView(PDFToolAPIView):
    """Generate page previews for PDF files."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            import base64
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            total_pages = len(doc)
            
            previews = []
            for page_num in range(total_pages):
                page = doc.load_page(page_num)
                
                # Render page to image at higher quality for sharper previews
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))  # 2.0x scaling for high quality
                img_bytes = pix.tobytes("png")  # Use PNG for lossless quality

                # Convert to base64
                img_base64 = base64.b64encode(img_bytes).decode('utf-8')

                previews.append({
                    'pageNumber': page_num + 1,
                    'image': f'data:image/png;base64,{img_base64}',
                    'width': pix.width,
                    'height': pix.height
                })
            
            doc.close()
            
            return Response({
                'previews': previews,
                'totalPages': total_pages
            })
            
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class OfficeFilePreviewView(PDFToolAPIView):
    """Generate first page preview for Office files (Word, Excel, PowerPoint)."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            import base64
            import traceback
            from apps.tools.converters.office_converter import convert_powerpoint_to_pdf, convert_excel_to_pdf
            from apps.tools.converters.word_to_pdf import convert_word_to_pdf
            
            # Determine file type and convert to PDF
            filename = file.name.lower()
            
            logger.info(f"Generating preview for: {filename}")
            
            if filename.endswith(('.doc', '.docx')):
                file.seek(0)
                pdf_bytes = convert_word_to_pdf(file)
            elif filename.endswith(('.xls', '.xlsx')):
                file.seek(0)
                pdf_bytes = convert_excel_to_pdf(file)
            elif filename.endswith(('.ppt', '.pptx')):
                logger.info("Converting PowerPoint to PDF for preview...")
                file.seek(0)
                try:
                    pdf_bytes = convert_powerpoint_to_pdf(file)
                    logger.info(f"PowerPoint converted, PDF size: {len(pdf_bytes)} bytes")
                except Exception as primary_error:
                    logger.warning(f"Primary PPT conversion failed, using fallback: {primary_error}")
                    try:
                        file.seek(0)
                        from apps.tools.converters.python_office_converter import convert_powerpoint_to_pdf_python
                        pdf_bytes = convert_powerpoint_to_pdf_python(file)
                        logger.info("Fallback PPT conversion succeeded")
                    except Exception as fallback_error:
                        raise Exception(f"PowerPoint preview failed: {fallback_error}")
            else:
                return Response({'error': 'Unsupported file type'}, status=status.HTTP_400_BAD_REQUEST)
            
            # Open the generated PDF and extract first page as image
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            
            if len(doc) == 0:
                return Response({'error': 'Generated PDF has no pages'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            
            logger.info(f"PDF opened successfully, pages: {len(doc)}")
            
            # Get first page
            page = doc.load_page(0)
            
            # Render page to image at high quality
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            img_bytes = pix.tobytes("png")
            
            # Convert to base64
            img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            
            preview_data = {
                'preview': f'data:image/png;base64,{img_base64}',
                'width': pix.width,
                'height': pix.height,
                'totalPages': len(doc)
            }
            
            doc.close()
            
            logger.info("Preview generated successfully")
            return Response(preview_data)
            
        except Exception as e:
            import traceback
            logger.error(f"Office file preview failed: {e}")
            logger.error(f"Traceback: {traceback.format_exc()}")
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class CompressPDFView(PDFToolAPIView):
    """Compress PDF to reduce file size."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from apps.tools.optimizers.compress import compress_pdf
            level = request.data.get('level', 'recommended')
            result = compress_pdf(file, level)
            
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_compressed.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class OrganizePDFView(PDFToolAPIView):
    """Organize PDF pages (reorder, rotate, delete)."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            import json
            
            pages_data = json.loads(request.data.get('pages', '[]'))
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            # Create new document with reorganized pages
            new_doc = fitz.open()
            for page_info in pages_data:
                page_num = page_info.get('page', 1) - 1
                rotation = page_info.get('rotation', 0)
                
                if 0 <= page_num < len(doc):
                    new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                    if rotation:
                        new_doc[-1].set_rotation(rotation)
            
            output = io.BytesIO()
            new_doc.save(output)
            new_doc.close()
            doc.close()
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="organized.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class FlattenPDFView(PDFToolAPIView):
    """Flatten PDF (merge annotations into content)."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            import fitz
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            for page in doc:
                # Flatten annotations
                for annot in page.annots():
                    annot.update()
            
            output = io.BytesIO()
            doc.save(output, deflate=True)
            doc.close()
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_flattened.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class CompressImageView(PDFToolAPIView):
    """Compress images."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        try:
            from PIL import Image
            
            level = request.data.get('level', 'recommended')
            quality_map = {'low': 30, 'recommended': 60, 'high': 85}
            quality = quality_map.get(level, 60)
            
            img = Image.open(file)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            
            output = io.BytesIO()
            img.save(output, format='JPEG', quality=quality, optimize=True)
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='image/jpeg')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_compressed.jpg"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# ─────────────────────────────────────────────────────────────────────────────
# PDF SECURITY
# ─────────────────────────────────────────────────────────────────────────────

class ProtectPDFView(PDFToolAPIView):
    """Add password protection to PDF."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        password = request.data.get('password')
        if not password:
            return Response({'error': 'Password required'}, status=status.HTTP_400_BAD_REQUEST)
        
        try:
            from apps.tools.security.protect import protect_pdf
            result = protect_pdf(file, password)
            
            response = HttpResponse(result, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_protected.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class UnlockPDFView(PDFToolAPIView):
    """Remove password protection from PDF."""
    
    def post(self, request):
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        password = request.data.get('password')
        if not password:
            return Response({'error': 'Password required'}, status=status.HTTP_400_BAD_REQUEST)
        
        try:
            import fitz
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            if doc.is_encrypted:
                if not doc.authenticate(password):
                    return Response({'error': 'Invalid password'}, status=status.HTTP_401_UNAUTHORIZED)
            
            output = io.BytesIO()
            doc.save(output)
            doc.close()
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_unlocked.pdf"'
            return response
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


# ─────────────────────────────────────────────────────────────────────────────
# AI/PREMIUM TOOLS
# ─────────────────────────────────────────────────────────────────────────────

class OCRPDFView(PDFToolAPIView):
    """Perform OCR on scanned PDF to make it searchable. Premium feature."""
    
    def post(self, request):
        # Premium check: require authenticated and premium user
        if not (request.user.is_authenticated and getattr(request.user, "is_premium", False)):
            return Response(
                {'error': 'OCR is a premium feature. Please upgrade your subscription.'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        language = request.data.get('language', 'eng')
        deskew = request.data.get('deskew', True)
        output_type = request.data.get('output', 'pdf')  # 'pdf' or 'text'
        
        try:
            from apps.tools.ai.ocr import ocr
            
            # Write to temp file since ocrmypdf needs file path
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_in:
                tmp_in.write(file.read())
                input_path = tmp_in.name
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
            
            result = ocr(input_path, output_path, language=language, deskew=deskew)
            
            # Clean up input file
            os.unlink(input_path)
            
            if not result.get('success'):
                if os.path.exists(output_path):
                    os.unlink(output_path)
                return Response(
                    {'error': result.get('message', 'OCR failed')},
                    status=status.HTTP_500_INTERNAL_SERVER_ERROR
                )
            
            if output_type == 'text':
                # Extract text from OCR'd PDF
                import fitz
                doc = fitz.open(output_path)
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()
                os.unlink(output_path)
                
                return Response({
                    'text': text,
                    'pages_processed': result.get('pages_processed'),
                    'language': language
                })
            else:
                # Return the OCR'd PDF
                with open(output_path, 'rb') as f:
                    pdf_content = f.read()
                os.unlink(output_path)
                
                response = HttpResponse(pdf_content, content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_ocr.pdf"'
                return response
            
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class EditPDFView(PDFToolAPIView):
    """Add text, highlights, and shapes to PDF. Premium feature."""
    
    def post(self, request):
        if not getattr(request.user, 'is_premium', False):
            return Response(
                {'error': 'Edit PDF is a premium feature. Please upgrade your subscription.'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        annotations = request.data.get('annotations', [])
        
        if not annotations:
            return Response({'error': 'No annotations provided'}, status=status.HTTP_400_BAD_REQUEST)
        
        try:
            import fitz
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            for annot in annotations:
                annot_type = annot.get('type')
                page_num = annot.get('page', 1) - 1  # Convert to 0-indexed
                
                if page_num < 0 or page_num >= len(doc):
                    continue
                
                page = doc[page_num]
                
                if annot_type == 'text':
                    # Add text annotation
                    x = annot.get('x', 0)
                    y = annot.get('y', 0)
                    content = annot.get('content', '')
                    font_size = annot.get('font_size', 12)
                    color = annot.get('color', '#000000')
                    
                    # Convert hex color to RGB tuple (0-1 range)
                    rgb = tuple(int(color.lstrip('#')[i:i+2], 16) / 255 for i in (0, 2, 4))
                    
                    page.insert_text(
                        fitz.Point(x, y),
                        content,
                        fontsize=font_size,
                        color=rgb
                    )
                
                elif annot_type == 'highlight':
                    rect = annot.get('rect', [0, 0, 100, 20])
                    page.add_highlight_annot(fitz.Rect(rect))
                
                elif annot_type == 'rectangle':
                    rect = annot.get('rect', [0, 0, 100, 100])
                    color = annot.get('color', '#FF0000')
                    rgb = tuple(int(color.lstrip('#')[i:i+2], 16) / 255 for i in (0, 2, 4))
                    
                    shape = page.new_shape()
                    shape.draw_rect(fitz.Rect(rect))
                    shape.finish(color=rgb, width=1)
                    shape.commit()
            
            output = io.BytesIO()
            doc.save(output)
            doc.close()
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_edited.pdf"'
            return response
            
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class RedactPDFView(PDFToolAPIView):
    """Redact sensitive information from PDF. Premium feature."""
    
    def post(self, request):
        if not request.user.is_premium:
            return Response(
                {'error': 'Redact PDF is a premium feature. Please upgrade your subscription.'},
                status=status.HTTP_403_FORBIDDEN
            )
        
        file, error = self.get_file_from_request(request)
        if error:
            return error
        
        redactions = request.data.get('redactions', [])
        
        if not redactions:
            return Response({'error': 'No redactions provided'}, status=status.HTTP_400_BAD_REQUEST)
        
        try:
            import fitz
            
            doc = fitz.open(stream=file.read(), filetype="pdf")
            
            for redact in redactions:
                page_num = redact.get('page', 1) - 1  # Convert to 0-indexed
                
                if page_num < 0 or page_num >= len(doc):
                    continue
                
                page = doc[page_num]
                
                if 'rect' in redact:
                    # Redact by coordinates
                    rect = fitz.Rect(redact['rect'])
                    page.add_redact_annot(rect, fill=(0, 0, 0))
                
                elif 'text' in redact:
                    # Redact by finding text
                    text_to_redact = redact['text']
                    text_instances = page.search_for(text_to_redact)
                    for inst in text_instances:
                        page.add_redact_annot(inst, fill=(0, 0, 0))
                
                # Apply redactions to this page
                page.apply_redactions()
            
            output = io.BytesIO()
            doc.save(output, garbage=4, deflate=True)  # Clean up and compress
            doc.close()
            output.seek(0)
            
            response = HttpResponse(output.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{file.name.rsplit(".", 1)[0]}_redacted.pdf"'
            return response
            
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


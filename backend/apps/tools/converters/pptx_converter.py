"""
Enhanced PowerPoint to PDF Converter using multiple strategies
"""
import io
import os
import sys
import logging
from PIL import Image
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import tempfile

logger = logging.getLogger(__name__)


def convert_pptx_to_pdf_via_images(file) -> bytes:
    """
    Convert PowerPoint to PDF by rendering slides as images.
    Optimized for speed - skips slow COM automation.
    
    Strategy:
    1. Extract images from PPTX zip structure (FAST)
    2. Convert images to PDF
    
    Args:
        file: Django UploadedFile or file-like object
    
    Returns:
        PDF bytes
    """
    content = file.read()

    # If the file is legacy PPT (not a zip), use Windows COM if available
    try:
        import zipfile
        if not zipfile.is_zipfile(io.BytesIO(content)):
            logger.info("Detected legacy PPT format; trying Windows COM conversion")
            if sys.platform == 'win32':
                try:
                    filename = getattr(file, 'name', 'presentation.ppt')
                    return _convert_via_windows_com(content, filename)
                except Exception as com_error:
                    logger.warning(f"Legacy PPT COM conversion failed: {com_error}")
                    raise Exception("Legacy PPT requires PowerPoint installed on Windows for previews.")
            else:
                raise Exception("Legacy PPT requires Windows PowerPoint for previews.")
    except Exception:
        pass
    
    # Skip slow COM automation - go directly to fast image extraction
    try:
        return _convert_via_image_extraction(content)
    except Exception as e:
        logger.error(f"Image extraction conversion failed: {e}")
        raise Exception(f"PowerPoint to PDF conversion failed: {str(e)}")


def _convert_via_windows_com(content: bytes, filename: str) -> bytes:
    """
    Convert using Windows COM automation (PowerPoint application).
    Only works on Windows with PowerPoint installed.
    """
    try:
        import comtypes.client
    except ImportError:
        raise Exception("comtypes not available")
    
    with tempfile.TemporaryDirectory() as tmpdir:
        # Save PPTX
        pptx_path = os.path.join(tmpdir, filename)
        with open(pptx_path, 'wb') as f:
            f.write(content)
        
        # Convert using PowerPoint
        pdf_path = os.path.join(tmpdir, 'output.pdf')
        
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        
        try:
            presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
            presentation.SaveAs(pdf_path, 32)  # 32 = ppSaveAsPDF
            presentation.Close()
        finally:
            powerpoint.Quit()
        
        # Read PDF
        with open(pdf_path, 'rb') as f:
            return f.read()


def _convert_via_image_extraction(content: bytes) -> bytes:
    """
    Convert PPTX to PDF by rendering slides.
    This creates a proper visual representation of each slide.
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import zipfile
    
    logger.info("Starting PPTX conversion...")
    
    try:
        prs = Presentation(io.BytesIO(content))
        logger.info(f"Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        logger.error(f"Failed to load presentation: {e}")
        raise
    
    # Create PDF
    pdf_buffer = io.BytesIO()
    
    # Use presentation dimensions if available, otherwise use standard 16:9
    try:
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height
        # Convert EMU to pixels (1 inch = 914400 EMU, 96 DPI)
        slide_width_px = int((slide_width_emu / 914400) * 96 * 2.5)  # 2.5x for better quality
        slide_height_px = int((slide_height_emu / 914400) * 96 * 2.5)
        logger.info(f"Slide dimensions: {slide_width_px}x{slide_height_px} px")
    except Exception as e:
        # Default to 16:9 at high resolution
        slide_width_px = 1920
        slide_height_px = 1080
        logger.info(f"Using default dimensions: {slide_width_px}x{slide_height_px} px")
    
    c = canvas.Canvas(pdf_buffer, pagesize=(slide_width_px * 0.75, slide_height_px * 0.75))  # Convert px to points
    
    slide_images = []
    
    # First try to extract any embedded slide thumbnails or preview images
    logger.info("Checking for embedded thumbnails...")
    embedded_images = _extract_slide_thumbnails(content)
    
    if embedded_images and len(embedded_images) >= len(prs.slides):
        # Use embedded thumbnails if available
        logger.info(f"Using {len(embedded_images)} embedded thumbnails")
        slide_images = embedded_images[:len(prs.slides)]
    else:
        # Render slides using content
        logger.info("Rendering slides from content...")
        for slide_idx, slide in enumerate(prs.slides):
            try:
                logger.info(f"Rendering slide {slide_idx + 1}/{len(prs.slides)}...")
                slide_img = _render_slide_with_content(slide, slide_width_px, slide_height_px)
                slide_images.append(slide_img)
                logger.info(f"Slide {slide_idx + 1} rendered successfully")
            except Exception as e:
                logger.error(f"Could not render slide {slide_idx + 1}: {e}")
                import traceback
                logger.error(traceback.format_exc())
                # Create a fallback slide
                fallback = Image.new('RGB', (slide_width_px, slide_height_px), 'white')
                draw = ImageDraw.Draw(fallback)
                try:
                    font = ImageFont.truetype("arial.ttf", 48)
                except:
                    font = ImageFont.load_default()
                draw.text((100, slide_height_px // 2), f"Slide {slide_idx + 1}", fill='black', font=font)
                slide_images.append(fallback)
    
    if not slide_images:
        logger.error("No slides were converted!")
        raise Exception("No slides could be converted")
    
    logger.info(f"Adding {len(slide_images)} slides to PDF...")
    
    # Add slides to PDF
    for idx, slide_img in enumerate(slide_images):
        try:
            # Convert PIL Image to bytes
            img_buffer = io.BytesIO()
            slide_img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            # Add to PDF
            img_reader = ImageReader(img_buffer)
            c.drawImage(img_reader, 0, 0, 
                       width=slide_width_px * 0.75, 
                       height=slide_height_px * 0.75)
            c.showPage()
        except Exception as e:
            logger.error(f"Failed to add slide {idx + 1} to PDF: {e}")
    
    c.save()
    pdf_bytes = pdf_buffer.getvalue()
    logger.info(f"PDF generated successfully, size: {len(pdf_bytes)} bytes")
    return pdf_bytes


def _render_slide_with_content(slide, width: int, height: int) -> Image:
    """
    Render a single slide with its text and shape content.
    Returns a PIL Image.
    """
    from PIL import Image, ImageDraw, ImageFont
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    # Create slide background
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    # Try to load fonts
    try:
        title_font = ImageFont.truetype("arial.ttf", int(height * 0.06))
        body_font = ImageFont.truetype("arial.ttf", int(height * 0.04))
        small_font = ImageFont.truetype("arial.ttf", int(height * 0.03))
    except:
        title_font = body_font = small_font = ImageFont.load_default()
    
    # Process shapes in order
    for shape in slide.shapes:
        try:
            # Get shape position and size (convert EMU to pixels)
            left = int((shape.left / 914400) * 96 * 2.5)
            top = int((shape.top / 914400) * 96 * 2.5)
            shape_width = int((shape.width / 914400) * 96 * 2.5)
            shape_height = int((shape.height / 914400) * 96 * 2.5)
            
            # Handle different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Try to get the image
                try:
                    image_data = shape.image.blob
                    shape_img = Image.open(io.BytesIO(image_data))
                    shape_img = shape_img.resize((shape_width, shape_height), Image.Resampling.LANCZOS)
                    img.paste(shape_img, (left, top))
                except:
                    # Draw placeholder for image
                    draw.rectangle([left, top, left + shape_width, top + shape_height], 
                                 outline='gray', width=2)
            
            elif hasattr(shape, "text") and shape.text.strip():
                # Text shape
                text = shape.text.strip()
                
                # Choose font based on text length and shape size
                if len(text) < 50 and shape_height > height * 0.15:
                    font = title_font
                elif len(text) < 100:
                    font = body_font
                else:
                    font = small_font
                
                # Word wrap text
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=font)
                    if bbox[2] - bbox[0] <= shape_width - 20:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw text lines
                y = top + 10
                line_height = int(height * 0.05)
                
                for line in lines[:10]:  # Limit to 10 lines per shape
                    if y + line_height > top + shape_height:
                        break
                    draw.text((left + 10, y), line, fill='black', font=font)
                    y += line_height
            
        except Exception as e:
            logger.debug(f"Could not process shape: {e}")
            continue
    
    return img


def _extract_slide_thumbnails(content: bytes) -> list:
    """
    Try to extract pre-rendered slide thumbnails from PPTX.
    Returns list of PIL Images.
    """
    import zipfile
    from PIL import Image
    
    images = []
    
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as pptx_zip:
            # Look for slide preview images in docProps/thumbnail.jpeg or ppt/slides/_rels/*.xml.rels
            thumbnail_files = []
            
            for file_info in pptx_zip.filelist:
                filename = file_info.filename.lower()
                # Look for thumbnail or preview images
                if 'thumbnail' in filename or (filename.startswith('ppt/slides/') and filename.endswith('.jpeg')):
                    thumbnail_files.append(file_info.filename)
            
            # Read thumbnail images
            for thumb_file in sorted(thumbnail_files):
                try:
                    img_data = pptx_zip.read(thumb_file)
                    img = Image.open(io.BytesIO(img_data))
                    images.append(img)
                except:
                    continue
    except Exception as e:
        logger.debug(f"Could not extract thumbnails: {e}")
    
    return images


def convert_powerpoint_to_pdf_best(file) -> bytes:
    """
    Main entry point - uses best available method.
    
    Priority:
    1. Image extraction from PPTX - Good quality
    2. Text-based rendering - Fallback
    
    Args:
        file: Django UploadedFile or file-like object
    
    Returns:
        PDF bytes
    """
    try:
        return convert_pptx_to_pdf_via_images(file)
    except Exception as e:
        logger.error(f"PowerPoint conversion failed with all methods: {e}")
        raise

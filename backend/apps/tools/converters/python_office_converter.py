"""
Pure Python Office Document Converters (Fallback when LibreOffice is not available)
"""
import io
import logging
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY

logger = logging.getLogger(__name__)


def convert_word_to_pdf_python(file) -> bytes:
    """
    Convert Word document to PDF using pure Python (docx2python + reportlab).
    
    Args:
        file: Django UploadedFile or file-like object
    
    Returns:
        PDF bytes
    """
    try:
        from docx import Document
        
        content = file.read()
        doc = Document(io.BytesIO(content))
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Add content
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Determine style based on paragraph properties
                style = styles['Normal']
                if paragraph.style.name.startswith('Heading'):
                    style = styles['Heading1']
                
                para = Paragraph(paragraph.text, style)
                story.append(para)
                story.append(Spacer(1, 0.2*inch))
        
        # Add tables
        for table in doc.tables:
            data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                data.append(row_data)
            
            if data:
                t = Table(data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(t)
                story.append(Spacer(1, 0.3*inch))
        
        pdf_doc.build(story)
        return pdf_buffer.getvalue()
        
    except Exception as e:
        logger.error(f"Pure Python Word to PDF failed: {e}")
        raise


def convert_excel_to_pdf_python(file) -> bytes:
    """
    Convert Excel spreadsheet to PDF using pure Python (openpyxl + reportlab).
    
    Args:
        file: Django UploadedFile or file-like object
    
    Returns:
        PDF bytes
    """
    try:
        from openpyxl import load_workbook
        
        content = file.read()
        wb = load_workbook(io.BytesIO(content), data_only=True)
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        # Process each sheet
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # Add sheet name as heading
            story.append(Paragraph(f"<b>{sheet_name}</b>", styles['Heading1']))
            story.append(Spacer(1, 0.3*inch))
            
            # Convert sheet to table data
            data = []
            for row in sheet.iter_rows(values_only=True):
                # Convert None to empty string and format values
                row_data = [str(cell) if cell is not None else '' for cell in row]
                # Skip completely empty rows
                if any(cell.strip() for cell in row_data):
                    data.append(row_data)
            
            if data:
                # Create table with styling
                t = Table(data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#136dec')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                    ('TOPPADDING', (0, 0), (-1, 0), 8),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ]))
                story.append(t)
            
            # Add page break between sheets
            if sheet_name != wb.sheetnames[-1]:
                story.append(PageBreak())
        
        pdf_doc.build(story)
        return pdf_buffer.getvalue()
        
    except Exception as e:
        logger.error(f"Pure Python Excel to PDF failed: {e}")
        raise


def convert_powerpoint_to_pdf_python(file) -> bytes:
    """
    Convert PowerPoint presentation to PDF using pure Python (python-pptx + reportlab).
    Enhanced version that properly extracts and displays all text content.
    
    Args:
        file: Django UploadedFile or file-like object
    
    Returns:
        PDF bytes
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        content = file.read()
        prs = Presentation(io.BytesIO(content))
        
        # Create PDF with better layout
        pdf_buffer = io.BytesIO()
        pdf_doc = SimpleDocTemplate(
            pdf_buffer, 
            pagesize=letter,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch,
            leftMargin=0.75*inch,
            rightMargin=0.75*inch
        )
        story = []
        styles = getSampleStyleSheet()
        
        # Create custom styles
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#136dec'),
            spaceAfter=20,
            spaceBefore=10,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'SlideHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#111418'),
            spaceAfter=12,
            bold=True
        )
        
        body_style = ParagraphStyle(
            'SlideBody',
            parent=styles['Normal'],
            fontSize=12,
            textColor=colors.HexColor('#333333'),
            spaceAfter=10,
            leading=16
        )
        
        # Process each slide
        for idx, slide in enumerate(prs.slides, 1):
            has_content = False
            
            # Check if slide has title
            if slide.shapes.title and slide.shapes.title.text.strip():
                # Add title
                title_text = slide.shapes.title.text.strip()
                story.append(Paragraph(title_text, title_style))
                story.append(Spacer(1, 0.3*inch))
                has_content = True
            else:
                # Add slide number as header if no title
                story.append(Paragraph(f"Slide {idx}", heading_style))
                story.append(Spacer(1, 0.2*inch))
            
            # Extract all text from shapes (excluding title)
            text_content = []
            tables_content = []
            
            for shape in slide.shapes:
                try:
                    # Skip title shape as we already processed it
                    if shape == slide.shapes.title:
                        continue
                    
                    # Handle text boxes and text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                # Check indentation level for bullets
                                level = paragraph.level
                                indent = "&nbsp;" * (level * 4)
                                
                                # Add bullet point
                                if level > 0:
                                    text_content.append(f"{indent}• {text}")
                                else:
                                    text_content.append(f"• {text}")
                                has_content = True
                    
                    # Handle tables
                    elif shape.has_table:
                        table = shape.table
                        data = []
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            if any(row_data):  # Only add non-empty rows
                                data.append(row_data)
                        
                        if data:
                            tables_content.append(data)
                            has_content = True
                    
                    # Handle grouped shapes
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for sub_shape in shape.shapes:
                            if sub_shape.has_text_frame:
                                text = sub_shape.text_frame.text.strip()
                                if text:
                                    text_content.append(f"• {text}")
                                    has_content = True
                
                except Exception as shape_error:
                    # Skip problematic shapes
                    logger.warning(f"Could not process shape: {shape_error}")
                    continue
            
            # Add all text content
            if text_content:
                for text in text_content:
                    try:
                        para = Paragraph(text, body_style)
                        story.append(para)
                    except Exception as para_error:
                        # Fallback for problematic text
                        safe_text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                        para = Paragraph(safe_text, body_style)
                        story.append(para)
                story.append(Spacer(1, 0.15*inch))
            
            # Add all tables
            for table_data in tables_content:
                try:
                    t = Table(table_data, hAlign='LEFT')
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#136dec')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                        ('TOPPADDING', (0, 0), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 0.2*inch))
                except Exception as table_error:
                    logger.warning(f"Could not render table: {table_error}")
            
            # If no content was extracted, add a note
            if not has_content:
                story.append(Paragraph(
                    f"<i>Slide {idx} contains visual elements that cannot be displayed in text format.</i>", 
                    body_style
                ))
                story.append(Spacer(1, 0.2*inch))
            
            # Add page break between slides (except for last slide)
            if idx < len(prs.slides):
                story.append(PageBreak())
        
        # Build PDF
        if not story:
            # Empty presentation
            story.append(Paragraph("This presentation appears to be empty or contains no extractable content.", body_style))
        
        pdf_doc.build(story)
        return pdf_buffer.getvalue()
        
    except Exception as e:
        logger.error(f"Pure Python PowerPoint to PDF failed: {e}")
        raise


import os
import fitz
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from io import BytesIO
import traceback

def create_dummy_pdf(filename):
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Hello World", fontsize=20, color=(0, 0, 1))
    doc.save(filename)
    doc.close()

def test_conversion():
    input_file = "debug_ppt_test.pdf"
    output_file = "debug_ppt_test.pptx"
    
    print(f"Creating dummy PDF: {input_file}")
    create_dummy_pdf(input_file)
    
    print("Starting conversion...")
    try:
        # Replicating the View Logic verbatim
        doc = fitz.open(input_file)
        prs = Presentation()
        
        if len(doc) > 0:
            page0 = doc[0]
            prs.slide_width = Pt(page0.rect.width)
            prs.slide_height = Pt(page0.rect.height)

        blank_layout = prs.slide_layouts[6] 

        for page in doc:
            slide = prs.slides.add_slide(blank_layout)

            page_dict = page.get_text("dict")
            
            for block in page_dict["blocks"]:
                bbox = block["bbox"]
                x, y, x1, y1 = bbox
                w_box = x1 - x
                h_box = y1 - y
                
                if block["type"] == 0: # TEXT
                    txBox = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w_box), Pt(h_box))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    for line in block["lines"]:
                        p = tf.add_paragraph()
                        for span in line["spans"]:
                            run = p.add_run()
                            run.text = span["text"]
                            run.font.size = Pt(span["size"])
                            run.font.name = span["font"]
                            
                            try:
                                c = span["color"]
                                r = (c >> 16) & 0xFF
                                g = (c >> 8) & 0xFF
                                b = c & 0xFF
                                run.font.color.rgb = RGBColor(r, g, b)
                            except:
                                pass
                                
                            flags = span["flags"]
                            if flags & 2**4: run.font.bold = True
                            if flags & 2**1: run.font.italic = True

                elif block["type"] == 1: # IMAGE
                    image_data = block["image"]
                    img_stream = BytesIO(image_data)
                    try:
                        slide.shapes.add_picture(img_stream, Pt(x), Pt(y), width=Pt(w_box), height=Pt(h_box))
                    except:
                        pass

        doc.close()
        prs.save(output_file)
        print(f"SUCCESS: Created {output_file}")

    except Exception as e:
        print("CRASH/ERROR DETAILS:")
        traceback.print_exc()

if __name__ == "__main__":
    test_conversion()

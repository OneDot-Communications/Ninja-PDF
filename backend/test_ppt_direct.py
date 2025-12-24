"""
Test PowerPoint COM directly
"""
import sys
import os
import tempfile

print("Testing PowerPoint COM automation...")

# Create a simple test PPTX first
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# Add title
from pptx.util import Inches
left = top = Inches(1)
width = Inches(8)
height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Test Slide - Hello World"

test_pptx = os.path.abspath("test_conversion.pptx")
test_pdf = os.path.abspath("test_conversion.pdf")

prs.save(test_pptx)
print(f"Created test PPTX: {test_pptx}")

# Now try COM conversion
print("\nAttempting PowerPoint COM conversion...")

try:
    import win32com.client
    import pythoncom
    
    # Initialize COM for this thread
    pythoncom.CoInitialize()
    
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    # Don't set Visible=False immediately, let it open naturally
    
    # Open presentation
    presentation = powerpoint.Presentations.Open(test_pptx, WithWindow=False)
    print("Opened presentation successfully")
    
    # Export to PDF
    presentation.SaveAs(test_pdf, 32)  # 32 = ppSaveAsPDF
    print("SaveAs PDF called")
    
    presentation.Close()
    powerpoint.Quit()
    
    pythoncom.CoUninitialize()
    
    if os.path.exists(test_pdf):
        size = os.path.getsize(test_pdf)
        print(f"\n✓ SUCCESS! PDF created: {test_pdf} ({size} bytes)")
    else:
        print("\n✗ FAILED: PDF not created")

except Exception as e:
    print(f"\n✗ ERROR: {e}")
    import traceback
    traceback.print_exc()

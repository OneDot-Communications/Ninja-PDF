# pdf_to_powerpoint.py
"""
PDF to PowerPoint conversion module.
Converts PDF pages to PowerPoint slides using python-pptx.
Each PDF page becomes a full-slide image in the presentation.
"""

import os
import io
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def parse_page_range(page_range, total_pages):
    """
    Parse page range string and return list of 1-based page numbers.
    
    Args:
        page_range: String like "1,3-5,8" or "all"
        total_pages: Total number of pages in PDF
    
    Returns:
        list: List of page numbers (1-based)
    """
    if not page_range or page_range.lower() == 'all':
        return list(range(1, total_pages + 1))
    
    pages = set()
    parts = page_range.split(',')
    
    for part in parts:
        part = part.strip()
        if '-' in part:
            try:
                start, end = map(int, part.split('-'))
                start = max(1, start)
                end = min(total_pages, end)
                for i in range(start, end + 1):
                    pages.add(i)
            except ValueError:
                continue
        else:
            try:
                page = int(part)
                if 1 <= page <= total_pages:
                    pages.add(page)
            except ValueError:
                continue
    
    return sorted(list(pages))


def pdf_page_to_image_bytes(pdf_doc, page_num, dpi=200):
    """
    Convert a PDF page to PNG image bytes.
    
    Args:
        pdf_doc: PyMuPDF document object
        page_num: Page number (1-based)
        dpi: DPI for rendering (higher = better quality)
    
    Returns:
        bytes: PNG image data
    """
    page = pdf_doc[page_num - 1]  # Convert to 0-based
    
    # Calculate zoom factor for desired DPI
    zoom = dpi / 72.0  # PDF default is 72 DPI
    mat = fitz.Matrix(zoom, zoom)
    
    # Render page to pixmap
    pix = page.get_pixmap(matrix=mat, alpha=False)
    
    # Convert to PNG bytes
    img_bytes = pix.tobytes("png")
    
    return img_bytes


def create_powerpoint_from_images(image_data_list, slide_width=10, slide_height=7.5):
    """
    Create a PowerPoint presentation from a list of image data.
    
    Args:
        image_data_list: List of tuples (page_num, image_bytes)
        slide_width: Slide width in inches (default: 10 inches - standard 4:3)
        slide_height: Slide height in inches (default: 7.5 inches - standard 4:3)
    
    Returns:
        bytes: PowerPoint file bytes
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    # Add each image as a full slide
    for page_num, img_bytes in image_data_list:
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]  # 6 = blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Open image to get dimensions
        img = Image.open(io.BytesIO(img_bytes))
        img_width, img_height = img.size
        
        # Calculate dimensions to fit slide while maintaining aspect ratio
        slide_width_px = prs.slide_width
        slide_height_px = prs.slide_height
        
        # Calculate scaling to fill the slide
        width_scale = slide_width_px / img_width
        height_scale = slide_height_px / img_height
        scale = min(width_scale, height_scale)
        
        # Calculate final dimensions
        final_width = int(img_width * scale)
        final_height = int(img_height * scale)
        
        # Center the image on the slide
        left = int((slide_width_px - final_width) / 2)
        top = int((slide_height_px - final_height) / 2)
        
        # Add image to slide
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, left, top, final_width, final_height)
    
    # Save to bytes
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer.read()


def convert_pdf_to_powerpoint(pdf_file, options=None):
    """
    Convert PDF file to PowerPoint presentation.
    
    Args:
        pdf_file: Django UploadedFile or file-like object
        options: Dictionary of conversion options
            - page_range: "1-5, 8" or "all" (default: "all")
            - format: "pptx" or "ppt" (default: "pptx")
            - dpi: Image quality in DPI (default: 200)
    
    Returns:
        tuple: (filename, pptx_bytes)
    """
    if options is None:
        options = {}
    
    # Get options
    page_range = options.get('page_range', 'all')
    output_format = options.get('format', 'pptx')
    dpi = int(options.get('dpi', 200))
    
    # Open PDF
    if hasattr(pdf_file, 'read'):
        pdf_file.seek(0)
        pdf_bytes = pdf_file.read()
        pdf_file.seek(0)
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    else:
        pdf_doc = fitz.open(pdf_file)
    
    try:
        total_pages = pdf_doc.page_count
        
        # Parse page range
        pages_to_convert = parse_page_range(page_range, total_pages)
        
        if not pages_to_convert:
            raise ValueError("No valid pages selected")
        
        # Convert each page to image
        image_data_list = []
        for page_num in pages_to_convert:
            img_bytes = pdf_page_to_image_bytes(pdf_doc, page_num, dpi=dpi)
            image_data_list.append((page_num, img_bytes))
        
        # Create PowerPoint
        pptx_bytes = create_powerpoint_from_images(image_data_list)
        
        # Generate filename
        original_name = getattr(pdf_file, 'name', 'document.pdf')
        filename = os.path.splitext(original_name)[0] + '.pptx'
        
        return filename, pptx_bytes
    
    finally:
        pdf_doc.close()
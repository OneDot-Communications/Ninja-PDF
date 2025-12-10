import os
import io
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF

def extract_slide_as_image(slide, slide_width, slide_height, dpi=150):
    """
    Extract slide content and render as an image.
    Uses high-resolution rendering for better quality.
    """
    # Calculate image dimensions based on DPI
    # Standard PowerPoint slide is 10" x 7.5"
    width_px = int((slide_width / 914400) * dpi * 10)
    height_px = int((slide_height / 914400) * dpi * 7.5)
    
    # Create a white background image
    img = Image.new('RGB', (width_px, height_px), 'white')
    draw = ImageDraw.Draw(img)
    
    # Try to load a font for text rendering
    try:
        font = ImageFont.truetype("arial.ttf", int(dpi / 5))
        title_font = ImageFont.truetype("arial.ttf", int(dpi / 3))
    except:
        font = ImageFont.load_default()
        title_font = ImageFont.load_default()
    
    # Extract and render slide content
    y_position = height_px // 10
    
    # Process each shape in the slide
    for shape in slide.shapes:
        try:
            # Handle text boxes and titles
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Determine if it's a title (usually larger and at top)
                is_title = hasattr(shape, "top") and shape.top < slide_height / 4
                current_font = title_font if is_title else font
                
                # Calculate text position (relative to slide dimensions)
                if hasattr(shape, "left") and hasattr(shape, "top"):
                    x = int((shape.left / slide_width) * width_px)
                    y = int((shape.top / slide_height) * height_px)
                else:
                    x = width_px // 20
                    y = y_position
                
                # Draw text with word wrapping
                max_width = width_px - (width_px // 10)
                lines = []
                words = text.split()
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=current_font)
                    if bbox[2] - bbox[0] <= max_width - x:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                            current_line = [word]
                        else:
                            lines.append(word)
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw each line
                for line in lines:
                    draw.text((x, y), line, fill='black', font=current_font)
                    bbox = draw.textbbox((0, 0), line, font=current_font)
                    y += (bbox[3] - bbox[1]) + 10
                
                y_position = y + 20
            
            # Handle images embedded in slides
            elif hasattr(shape, "image"):
                try:
                    image_stream = io.BytesIO(shape.image.blob)
                    slide_image = Image.open(image_stream)
                    
                    # Calculate image position and size
                    if hasattr(shape, "left") and hasattr(shape, "top"):
                        img_x = int((shape.left / slide_width) * width_px)
                        img_y = int((shape.top / slide_height) * height_px)
                        img_w = int((shape.width / slide_width) * width_px)
                        img_h = int((shape.height / slide_height) * height_px)
                        
                        # Resize and paste the image
                        slide_image = slide_image.resize((img_w, img_h), Image.Resampling.LANCZOS)
                        img.paste(slide_image, (img_x, img_y))
                except Exception as e:
                    print(f"Warning: Could not process image in shape: {e}")
        
        except Exception as e:
            print(f"Warning: Could not process shape: {e}")
            continue
    
    return img

def convert_powerpoint_to_pdf(input_path, output_path):
    """
    Convert PowerPoint (.pptx) to PDF using pure Python.
    Renders each slide as a high-quality image and combines them into a PDF.
    No external dependencies required (LibreOffice, MS Office, etc.)
    """
    
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"File not found: {input_path}")
    
    ext = os.path.splitext(input_path)[1].lower()
    if ext not in ['.pptx', '.ppt']:
        raise ValueError(f"Only .pptx and .ppt files are supported. Got: {ext}")
    
    # Create output directory if needed
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    try:
        print(f"Loading PowerPoint file: {input_path}")
        
        # Load the presentation
        prs = Presentation(input_path)
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        print(f"Presentation loaded: {len(prs.slides)} slides")
        print(f"Slide dimensions: {slide_width} x {slide_height}")
        
        # Create PDF document
        pdf_doc = fitz.open()
        
        # Process each slide
        for idx, slide in enumerate(prs.slides):
            print(f"Processing slide {idx + 1}/{len(prs.slides)}")
            
            # Render slide as image
            img = extract_slide_as_image(slide, slide_width, slide_height, dpi=150)
            
            # Convert PIL Image to bytes
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG', optimize=True, quality=95)
            img_bytes.seek(0)
            img_data = img_bytes.read()
            
            # Create a new PDF page with the image dimensions
            # Get image dimensions
            img_width, img_height = img.size
            
            # Create a new page in the PDF with the correct dimensions
            page = pdf_doc.new_page(width=img_width, height=img_height)
            
            # Insert the image into the page
            page.insert_image(
                fitz.Rect(0, 0, img_width, img_height),
                stream=img_data,
                keep_proportion=True
            )
        
        # Save the PDF
        pdf_doc.save(output_path, garbage=4, deflate=True)
        pdf_doc.close()
        
        print(f"Conversion successful: {output_path}")
        print(f"Created {len(prs.slides)}-page PDF from PowerPoint")
        
        return output_path
        
    except Exception as e:
        raise RuntimeError(f"PowerPoint to PDF conversion failed: {str(e)}")
